"""
H-001 Automated Insight Engine
A comprehensive AdTech data analysis and reporting system
"""

import os
import sys
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import warnings
from scipy import stats
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import IsolationForest
import json

# Report generation libraries
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Gemini API
import google.generativeai as genai

warnings.filterwarnings('ignore')

# Configuration
OUTPUT_DIR = "outputs"
PLOTS_DIR = os.path.join(OUTPUT_DIR, "plots")
TABLES_DIR = os.path.join(OUTPUT_DIR, "tables")

# Create necessary directories
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(PLOTS_DIR, exist_ok=True)
os.makedirs(TABLES_DIR, exist_ok=True)


class AdTechInsightEngine:
    """Main engine for automated AdTech data analysis and reporting"""
    
    def __init__(self, csv_path, gemini_api_key=None):
        self.csv_path = csv_path
        self.df = None
        self.analysis_results = {}
        self.insights = []
        self.gemini_api_key = gemini_api_key
        
        if gemini_api_key:
            genai.configure(api_key=gemini_api_key)
            self.model = genai.GenerativeModel('gemini-pro')
        else:
            self.model = None
        
        print(f"🚀 Initialized AdTech Insight Engine")
        print(f"📁 Output directory: {OUTPUT_DIR}")
    
    def load_csv(self):
        """Load and validate CSV data"""
        print(f"\n📊 Loading data from {self.csv_path}...")
        try:
            self.df = pd.read_csv(self.csv_path)
            print(f"✅ Loaded {len(self.df)} rows and {len(self.df.columns)} columns")
            print(f"📋 Columns: {', '.join(self.df.columns.tolist())}")
            
            # Attempt to parse date columns
            date_columns = [col for col in self.df.columns if 'date' in col.lower()]
            for col in date_columns:
                try:
                    self.df[col] = pd.to_datetime(self.df[col])
                    print(f"📅 Parsed date column: {col}")
                except:
                    pass
            
            return True
        except Exception as e:
            print(f"❌ Error loading CSV: {e}")
            return False
    
    def generate_summary_stats(self):
        """Generate comprehensive summary statistics"""
        print("\n📈 Generating summary statistics...")
        
        stats_dict = {
            'total_rows': len(self.df),
            'total_columns': len(self.df.columns),
            'numeric_columns': len(self.df.select_dtypes(include=[np.number]).columns),
            'categorical_columns': len(self.df.select_dtypes(include=['object']).columns),
            'missing_values': self.df.isnull().sum().sum(),
            'memory_usage_mb': self.df.memory_usage(deep=True).sum() / 1024**2
        }
        
        # Numeric summary
        numeric_df = self.df.select_dtypes(include=[np.number])
        if not numeric_df.empty:
            stats_dict['numeric_summary'] = numeric_df.describe().to_dict()
        
        self.analysis_results['summary_stats'] = stats_dict
        
        # Save summary table
        summary_df = pd.DataFrame([stats_dict])
        summary_df.to_csv(os.path.join(TABLES_DIR, 'summary_stats.csv'), index=False)
        
        print(f"✅ Summary statistics generated")
        return stats_dict
    
    def calculate_derived_metrics(self):
        """Calculate marketing metrics like CTR, CPC, CPM, ROAS, etc."""
        print("\n🔢 Calculating derived metrics...")
        
        df = self.df.copy()
        metrics_calculated = []
        
        # Common column name variations
        col_map = {col.lower().replace(' ', '_'): col for col in df.columns}
        
        # CTR (Click-Through Rate)
        if 'clicks' in col_map and 'impressions' in col_map:
            clicks_col = col_map['clicks']
            impressions_col = col_map['impressions']
            df['CTR'] = (df[clicks_col] / df[impressions_col] * 100).round(2)
            metrics_calculated.append('CTR (%)')
        
        # CPC (Cost Per Click)
        if 'cost' in col_map or 'spend' in col_map:
            cost_col = col_map.get('cost') or col_map.get('spend')
            if 'clicks' in col_map:
                clicks_col = col_map['clicks']
                df['CPC'] = (df[cost_col] / df[clicks_col]).round(2)
                metrics_calculated.append('CPC')
        
        # CPM (Cost Per Mille/Thousand Impressions)
        if ('cost' in col_map or 'spend' in col_map) and 'impressions' in col_map:
            cost_col = col_map.get('cost') or col_map.get('spend')
            impressions_col = col_map['impressions']
            df['CPM'] = (df[cost_col] / df[impressions_col] * 1000).round(2)
            metrics_calculated.append('CPM')
        
        # Conversion Rate
        if 'conversions' in col_map and 'clicks' in col_map:
            conversions_col = col_map['conversions']
            clicks_col = col_map['clicks']
            df['Conversion_Rate'] = (df[conversions_col] / df[clicks_col] * 100).round(2)
            metrics_calculated.append('Conversion Rate (%)')
        
        # ROAS (Return on Ad Spend)
        if 'revenue' in col_map and ('cost' in col_map or 'spend' in col_map):
            revenue_col = col_map['revenue']
            cost_col = col_map.get('cost') or col_map.get('spend')
            df['ROAS'] = (df[revenue_col] / df[cost_col]).round(2)
            metrics_calculated.append('ROAS')
        
        # CPA (Cost Per Acquisition)
        if ('cost' in col_map or 'spend' in col_map) and 'conversions' in col_map:
            cost_col = col_map.get('cost') or col_map.get('spend')
            conversions_col = col_map['conversions']
            df['CPA'] = (df[cost_col] / df[conversions_col]).round(2)
            metrics_calculated.append('CPA')
        
        self.df = df
        self.analysis_results['derived_metrics'] = metrics_calculated
        
        if metrics_calculated:
            print(f"✅ Calculated metrics: {', '.join(metrics_calculated)}")
        else:
            print("⚠️  No standard metric columns found")
        
        return metrics_calculated
    
    def generate_correlations(self):
        """Generate correlation matrix for numeric columns"""
        print("\n🔗 Generating correlation analysis...")
        
        numeric_df = self.df.select_dtypes(include=[np.number])
        
        if numeric_df.shape[1] < 2:
            print("⚠️  Not enough numeric columns for correlation")
            return None
        
        corr_matrix = numeric_df.corr()
        
        # Plot correlation heatmap
        plt.figure(figsize=(12, 10))
        sns.heatmap(corr_matrix, annot=True, fmt='.2f', cmap='coolwarm', 
                    center=0, square=True, linewidths=1)
        plt.title('Correlation Matrix - AdTech Metrics', fontsize=16, fontweight='bold')
        plt.tight_layout()
        plt.savefig(os.path.join(PLOTS_DIR, 'correlation_matrix.png'), dpi=300, bbox_inches='tight')
        plt.close()
        
        # Save correlation table
        corr_matrix.to_csv(os.path.join(TABLES_DIR, 'correlation_matrix.csv'))
        
        self.analysis_results['correlation_matrix'] = corr_matrix
        print("✅ Correlation analysis complete")
        
        return corr_matrix
    
    def detect_outliers(self):
        """Detect outliers using IQR and Isolation Forest methods"""
        print("\n🎯 Detecting outliers...")
        
        numeric_df = self.df.select_dtypes(include=[np.number])
        
        if numeric_df.empty:
            print("⚠️  No numeric columns for outlier detection")
            return None
        
        outliers_summary = {}
        
        # IQR Method
        for col in numeric_df.columns:
            Q1 = numeric_df[col].quantile(0.25)
            Q3 = numeric_df[col].quantile(0.75)
            IQR = Q3 - Q1
            lower_bound = Q1 - 1.5 * IQR
            upper_bound = Q3 + 1.5 * IQR
            
            outliers = numeric_df[(numeric_df[col] < lower_bound) | (numeric_df[col] > upper_bound)]
            outliers_summary[col] = {
                'count': len(outliers),
                'percentage': (len(outliers) / len(numeric_df) * 100),
                'lower_bound': lower_bound,
                'upper_bound': upper_bound
            }
        
        # Isolation Forest (multivariate)
        if numeric_df.shape[1] >= 2:
            iso_forest = IsolationForest(contamination=0.1, random_state=42)
            outlier_labels = iso_forest.fit_predict(numeric_df.fillna(0))
            outliers_summary['isolation_forest_outliers'] = (outlier_labels == -1).sum()
        
        self.analysis_results['outliers'] = outliers_summary
        
        # Save outlier summary
        outlier_df = pd.DataFrame(outliers_summary).T
        outlier_df.to_csv(os.path.join(TABLES_DIR, 'outlier_summary.csv'))
        
        print(f"✅ Outlier detection complete")
        return outliers_summary
    
    def perform_aggregations(self):
        """Perform aggregations by campaigns, ad groups, categories"""
        print("\n📊 Performing aggregations...")
        
        # Identify grouping columns
        group_cols = [col for col in self.df.columns if any(keyword in col.lower() 
                      for keyword in ['campaign', 'ad_group', 'category', 'source', 'channel'])]
        
        numeric_cols = self.df.select_dtypes(include=[np.number]).columns.tolist()
        
        aggregations = {}
        
        if group_cols and numeric_cols:
            for group_col in group_cols[:3]:  # Limit to top 3 grouping columns
                try:
                    agg_df = self.df.groupby(group_col)[numeric_cols].sum().round(2)
                    agg_df = agg_df.sort_values(by=numeric_cols[0], ascending=False).head(20)
                    
                    # Save aggregation
                    filename = f"aggregation_{group_col.lower().replace(' ', '_')}.csv"
                    agg_df.to_csv(os.path.join(TABLES_DIR, filename))
                    aggregations[group_col] = agg_df
                    
                    print(f"  ✓ Aggregated by {group_col}")
                except Exception as e:
                    print(f"  ⚠️  Could not aggregate by {group_col}: {e}")
        
        self.analysis_results['aggregations'] = aggregations
        return aggregations
    
    def analyze_time_series(self):
        """Analyze time-series trends if date columns exist"""
        print("\n📅 Analyzing time-series trends...")
        
        date_cols = self.df.select_dtypes(include=['datetime64']).columns.tolist()
        
        if not date_cols:
            print("⚠️  No date columns found for time-series analysis")
            return None
        
        date_col = date_cols[0]
        numeric_cols = self.df.select_dtypes(include=[np.number]).columns.tolist()[:5]
        
        if not numeric_cols:
            print("⚠️  No numeric columns for time-series")
            return None
        
        # Sort by date
        df_sorted = self.df.sort_values(by=date_col)
        
        # Plot time series
        fig, axes = plt.subplots(len(numeric_cols), 1, figsize=(14, 4*len(numeric_cols)))
        if len(numeric_cols) == 1:
            axes = [axes]
        
        for idx, col in enumerate(numeric_cols):
            axes[idx].plot(df_sorted[date_col], df_sorted[col], linewidth=2)
            axes[idx].set_title(f'{col} Over Time', fontweight='bold')
            axes[idx].set_xlabel('Date')
            axes[idx].set_ylabel(col)
            axes[idx].grid(True, alpha=0.3)
        
        plt.tight_layout()
        plt.savefig(os.path.join(PLOTS_DIR, 'time_series_trends.png'), dpi=300, bbox_inches='tight')
        plt.close()
        
        print("✅ Time-series analysis complete")
        return True
    
    def identify_top_performers(self):
        """Identify top and bottom performing segments"""
        print("\n🏆 Identifying top/bottom performers...")
        
        numeric_cols = self.df.select_dtypes(include=[np.number]).columns.tolist()
        
        if not numeric_cols:
            print("⚠️  No numeric columns for performance analysis")
            return None
        
        performers = {}
        
        # For each numeric metric, find top and bottom performers
        for metric in numeric_cols[:5]:  # Limit to 5 metrics
            sorted_df = self.df.sort_values(by=metric, ascending=False)
            
            performers[metric] = {
                'top_5': sorted_df.head(5)[[metric]].to_dict(),
                'bottom_5': sorted_df.tail(5)[[metric]].to_dict()
            }
        
        self.analysis_results['performers'] = performers
        print("✅ Performance analysis complete")
        
        return performers
    
    def generate_graphs(self):
        """Generate comprehensive visualizations"""
        print("\n📊 Generating visualizations...")
        
        numeric_cols = self.df.select_dtypes(include=[np.number]).columns.tolist()
        
        if not numeric_cols:
            print("⚠️  No numeric columns for visualization")
            return
        
        # 1. Distribution plots for top numeric columns
        fig, axes = plt.subplots(2, 2, figsize=(14, 10))
        axes = axes.flatten()
        
        for idx, col in enumerate(numeric_cols[:4]):
            axes[idx].hist(self.df[col].dropna(), bins=30, edgecolor='black')
            axes[idx].set_title(f'Distribution: {col}', fontweight='bold')
            axes[idx].set_xlabel(col)
            axes[idx].set_ylabel('Frequency')
            axes[idx].grid(True, alpha=0.3)
        
        plt.tight_layout()
        plt.savefig(os.path.join(PLOTS_DIR, 'distributions.png'), dpi=300, bbox_inches='tight')
        plt.close()
        
        # 2. Box plots for outlier visualization
        if len(numeric_cols) >= 2:
            fig, ax = plt.subplots(figsize=(14, 6))
            self.df[numeric_cols[:6]].boxplot(ax=ax)
            ax.set_title('Box Plot - Outlier Detection', fontsize=16, fontweight='bold')
            ax.set_ylabel('Values')
            plt.xticks(rotation=45, ha='right')
            plt.tight_layout()
            plt.savefig(os.path.join(PLOTS_DIR, 'boxplots.png'), dpi=300, bbox_inches='tight')
            plt.close()
        
        # 3. Top performers bar chart
        if len(numeric_cols) >= 1:
            top_metric = numeric_cols[0]
            top_10 = self.df.nlargest(10, top_metric)
            
            plt.figure(figsize=(12, 6))
            if 'campaign' in [col.lower() for col in self.df.columns]:
                campaign_col = [col for col in self.df.columns if col.lower() == 'campaign'][0]
                plt.barh(range(len(top_10)), top_10[top_metric])
                plt.yticks(range(len(top_10)), top_10[campaign_col])
            else:
                plt.barh(range(len(top_10)), top_10[top_metric])
                plt.yticks(range(len(top_10)), top_10.index)
            
            plt.xlabel(top_metric)
            plt.title(f'Top 10 by {top_metric}', fontsize=16, fontweight='bold')
            plt.tight_layout()
            plt.savefig(os.path.join(PLOTS_DIR, 'top_performers.png'), dpi=300, bbox_inches='tight')
            plt.close()
        
        print("✅ Visualizations generated")
    
    def generate_insights_with_ai(self):
        """Generate insights using Gemini AI"""
        print("\n🤖 Generating AI insights...")
        
        if not self.model:
            print("⚠️  Gemini API not configured, generating basic insights")
            return self.generate_basic_insights()
        
        # Prepare analysis summary for AI
        summary = f"""
        Analyze this AdTech dataset and provide executive-level insights:
        
        Dataset Overview:
        - Total Rows: {len(self.df)}
        - Total Columns: {len(self.df.columns)}
        - Columns: {', '.join(self.df.columns.tolist())}
        
        Summary Statistics:
        {self.df.describe().to_string()}
        
        Derived Metrics Calculated:
        {', '.join(self.analysis_results.get('derived_metrics', []))}
        
        Please provide:
        1. Key findings and trends
        2. Performance insights
        3. Areas of concern
        4. Actionable recommendations
        
        Format as clear bullet points.
        """
        
        try:
            response = self.model.generate_content(summary)
            ai_insights = response.text
            self.insights.append("AI-Generated Insights:")
            self.insights.append(ai_insights)
            print("✅ AI insights generated")
        except Exception as e:
            print(f"⚠️  AI generation failed: {e}, using basic insights")
            return self.generate_basic_insights()
        
        return self.insights
    
    def generate_basic_insights(self):
        """Generate basic insights without AI"""
        insights = []
        
        insights.append("📊 Key Findings:")
        insights.append(f"• Dataset contains {len(self.df)} records across {len(self.df.columns)} dimensions")
        
        numeric_cols = self.df.select_dtypes(include=[np.number]).columns.tolist()
        if numeric_cols:
            top_metric = numeric_cols[0]
            total_value = self.df[top_metric].sum()
            avg_value = self.df[top_metric].mean()
            insights.append(f"• Total {top_metric}: {total_value:,.2f}")
            insights.append(f"• Average {top_metric}: {avg_value:,.2f}")
        
        if 'derived_metrics' in self.analysis_results:
            insights.append(f"• Calculated {len(self.analysis_results['derived_metrics'])} derived metrics")
        
        self.insights = insights
        return insights
    
    def build_pdf_report(self):
        """Generate comprehensive PDF report"""
        print("\n📄 Building PDF report...")
        
        pdf_path = os.path.join(OUTPUT_DIR, 'AdTech_Insight_Report.pdf')
        doc = SimpleDocTemplate(pdf_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a1a1a'),
            spaceAfter=30,
            alignment=1
        )
        story.append(Paragraph("AdTech Insight Report", title_style))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", 
                               styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", styles['Heading2']))
        for insight in self.insights:
            story.append(Paragraph(insight, styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        
        story.append(PageBreak())
        
        # Summary Statistics
        story.append(Paragraph("Summary Statistics", styles['Heading2']))
        if 'summary_stats' in self.analysis_results:
            stats = self.analysis_results['summary_stats']
            for key, value in stats.items():
                if key != 'numeric_summary':
                    story.append(Paragraph(f"<b>{key}:</b> {value}", styles['Normal']))
        
        story.append(PageBreak())
        
        # Add visualizations
        story.append(Paragraph("Visual Analysis", styles['Heading2']))
        
        plot_files = [
            ('correlation_matrix.png', 'Correlation Matrix'),
            ('distributions.png', 'Metric Distributions'),
            ('boxplots.png', 'Outlier Analysis'),
            ('top_performers.png', 'Top Performers'),
            ('time_series_trends.png', 'Time Series Trends')
        ]
        
        for plot_file, caption in plot_files:
            plot_path = os.path.join(PLOTS_DIR, plot_file)
            if os.path.exists(plot_path):
                story.append(Paragraph(caption, styles['Heading3']))
                img = Image(plot_path, width=6*inch, height=4*inch)
                story.append(img)
                story.append(Spacer(1, 0.3*inch))
        
        # Build PDF
        doc.build(story)
        print(f"✅ PDF report saved: {pdf_path}")
        return pdf_path
    
    def build_ppt_report(self):
        """Generate PowerPoint presentation"""
        print("\n📊 Building PowerPoint presentation...")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "AdTech Insight Report"
        subtitle.text = f"Automated Analysis Report\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Executive Summary Slide
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = summary_slide.shapes.title
        title.text = "Executive Summary"
        
        content = summary_slide.placeholders[1]
        tf = content.text_frame
        tf.text = "\n".join(self.insights[:10])
        
        # Visualizations Slides
        plot_files = [
            ('correlation_matrix.png', 'Correlation Analysis'),
            ('distributions.png', 'Metric Distributions'),
            ('top_performers.png', 'Top Performers'),
            ('time_series_trends.png', 'Trends Over Time')
        ]
        
        for plot_file, slide_title in plot_files:
            plot_path = os.path.join(PLOTS_DIR, plot_file)
            if os.path.exists(plot_path):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = slide_title
                
                left = Inches(1)
                top = Inches(1.5)
                pic = slide.shapes.add_picture(plot_path, left, top, width=Inches(8))
        
        # Summary Statistics Slide
        stats_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = stats_slide.shapes.title
        title.text = "Key Metrics"
        
        if 'summary_stats' in self.analysis_results:
            stats = self.analysis_results['summary_stats']
            content = stats_slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Total Rows: {stats.get('total_rows', 'N/A')}\n"
            tf.text += f"Total Columns: {stats.get('total_columns', 'N/A')}\n"
            tf.text += f"Numeric Columns: {stats.get('numeric_columns', 'N/A')}\n"
            if 'derived_metrics' in self.analysis_results:
                tf.text += f"\nDerived Metrics: {', '.join(self.analysis_results['derived_metrics'])}"
        
        # Save presentation
        ppt_path = os.path.join(OUTPUT_DIR, 'AdTech_Insight_Presentation.pptx')
        prs.save(ppt_path)
        print(f"✅ PowerPoint saved: {ppt_path}")
        return ppt_path
    
    def run_full_analysis(self):
        """Execute complete analysis pipeline"""
        print("\n" + "="*60)
        print("🚀 STARTING AUTOMATED INSIGHT ENGINE")
        print("="*60)
        
        # Step 1: Load data
        if not self.load_csv():
            return False
        
        # Step 2: Generate summary statistics
        self.generate_summary_stats()
        
        # Step 3: Calculate derived metrics
        self.calculate_derived_metrics()
        
        # Step 4: Generate correlations
        self.generate_correlations()
        
        # Step 5: Detect outliers
        self.detect_outliers()
        
        # Step 6: Perform aggregations
        self.perform_aggregations()
        
        # Step 7: Time series analysis
        self.analyze_time_series()
        
        # Step 8: Identify top performers
        self.identify_top_performers()
        
        # Step 9: Generate visualizations
        self.generate_graphs()
        
        # Step 10: Generate insights
        self.generate_insights_with_ai()
        
        # Step 11: Build PDF report
        self.build_pdf_report()
        
        # Step 12: Build PowerPoint
        self.build_ppt_report()
        
        print("\n" + "="*60)
        print("✅ ANALYSIS COMPLETE!")
        print("="*60)
        print(f"\n📁 All outputs saved in: {OUTPUT_DIR}/")
        print(f"   - PDF Report: AdTech_Insight_Report.pdf")
        print(f"   - PowerPoint: AdTech_Insight_Presentation.pptx")
        print(f"   - Visualizations: plots/")
        print(f"   - Data Tables: tables/")
        
        return True


def main():
    """Main execution function"""
    print("""
    ╔═══════════════════════════════════════════════════════════╗
    ║         H-001 AUTOMATED INSIGHT ENGINE                    ║
    ║         AdTech Data Analysis & Reporting System           ║
    ╚═══════════════════════════════════════════════════════════╝
    """)
    
    # Get inputs
    if len(sys.argv) > 1:
        csv_path = sys.argv[1]
    else:
        csv_path = input("📁 Enter CSV file path: ").strip()
    
    gemini_key = os.getenv('GEMINI_API_KEY')
    if not gemini_key:
        gemini_key = input("🔑 Enter Gemini API key (optional, press Enter to skip): ").strip()
        if not gemini_key:
            gemini_key = None
    
    # Initialize and run engine
    engine = AdTechInsightEngine(csv_path, gemini_api_key=gemini_key)
    success = engine.run_full_analysis()
    
    if success:
        print("\n✨ Analysis completed successfully!")
    else:
        print("\n❌ Analysis failed. Please check your data and try again.")


if __name__ == "__main__":
    main()